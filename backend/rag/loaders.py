import os
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation

def load_pdf(file_path: str) -> list[Document]:
    """Loads a PDF document and returns normalized Document objects."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    loader = PyPDFLoader(file_path)
    docs = loader.load()
    
    normalized_docs = []
    for doc in docs:
        # PyPDFLoader returns 0-indexed page metadata, convert it to 1-indexed
        page_num = doc.metadata.get("page", 0) + 1
        normalized_docs.append(
            Document(
                page_content=doc.page_content,
                metadata={
                    "source": os.path.basename(file_path),
                    "page": page_num,
                    "type": "pdf"
                }
            )
        )
    return normalized_docs

def load_pptx(file_path: str) -> list[Document]:
    """Loads a PPT/PPTX document slide-by-slide and returns Document objects."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Failed to parse PPT/PPTX file: {str(e)}")
        
    documents = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text.strip())
        
        combined_text = "\n".join(slide_text_parts)
        if combined_text.strip():
            documents.append(
                Document(
                    page_content=combined_text,
                    metadata={
                        "source": os.path.basename(file_path),
                        "page": slide_idx + 1,  # 1-indexed slide number
                        "type": "pptx"
                    }
                )
            )
    return documents

def load_document(file_path: str) -> list[Document]:
    """Determines file type and loads it accordingly."""
    _, ext = os.path.splitext(file_path.lower())
    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext in [".ppt", ".pptx"]:
        return load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
